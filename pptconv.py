import streamlit as st
from pptx import Presentation
from pptx.util import Pt
from openai import OpenAI
import io
import os

# Load OpenAI API key from environment variable
if "openai_api_key" not in st.session_state:
    st.session_state.openai_api_key = os.environ.get("OPENAI_API_KEY")
client = OpenAI(api_key=st.session_state.openai_api_key)

def translate_text(text, target_language="English"):
    """
    GPT를 사용하여 텍스트를 번역합니다.
    """
    prompt = f"Translate the following Korean text to {target_language}:\n\n{text}"

    try:
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",  # 원하는 모델로 변경 가능
            messages=[
                {"role": "system", "content": "You are a helpful assistant that translates Korean to English."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=2000,
            n=1,
            stop=None,
            temperature=0.3,
        )
        translated_text = response.choices[0].message.content.strip()
        return translated_text
    except Exception as e:
        st.error(f"번역 중 오류 발생: {e}")
        return text  # 오류 발생 시 원본 텍스트 반환

def extract_text_from_ppt(ppt_stream):
    """
    PPT 파일에서 텍스트를 추출합니다.
    """
    prs = Presentation(ppt_stream)
    slides_text = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides_text.append(slide_text)
    return slides_text

def create_ppt_with_translated_text(original_ppt_stream, translated_slides_text):
    """
    번역된 텍스트로 새로운 PPT 파일을 생성합니다.
    """
    prs = Presentation(original_ppt_stream)

    for slide, translated_texts in zip(prs.slides, translated_slides_text):
        text_idx = 0
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if text_idx < len(translated_texts):
                    shape.text = translated_texts[text_idx]
                    # 글꼴 크기나 스타일 조정
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(10)
                    text_idx += 1
    # PPT을 메모리에 저장
    output_stream = io.BytesIO()
    prs.save(output_stream)
    output_stream.seek(0)
    return output_stream

def main():
    st.set_page_config(page_title="한글 PPT를 영어로 변환", layout="centered",page_icon="✨")
    st.title("PPT Translator for AGS")
    st.write("이일범 짱.")

    # 파일 업로드
    uploaded_file = st.file_uploader("한글 PPT 파일을 업로드하세요 (.pptx)", type=["pptx"])

    if uploaded_file is not None:
        if uploaded_file.name.split(".")[-1].lower() != "pptx":
            st.error("지원되지 않는 파일 형식입니다. .pptx 파일을 업로드해주세요.")
            return

        with st.spinner("PPT 파일에서 텍스트를 추출하는 중..."):
            try:
                slides_text = extract_text_from_ppt(uploaded_file)
            except Exception as e:
                st.error(f"PPT 파일을 처리하는 중 오류가 발생했습니다: {e}")
                return

        st.success("텍스트 추출 완료!")

        translated_slides = []
        for idx, slide_text in enumerate(slides_text):
            st.write(f"슬라이드 {idx+1} 번역 중...")
            translated_text = []
            for text in slide_text:
                if text.strip() == "":
                    translated_text.append("")
                    continue
                translated = translate_text(text)
                translated_text.append(translated)
            translated_slides.append(translated_text)

        # 재원 PPT 파일 스트리밍 다시 읽기 위해 파일 포인터를 처음으로 되돌림
        uploaded_file.seek(0)

        with st.spinner("번역된 텍스트로 새로운 PPT 파일을 생성하는 중..."):
            try:
                translated_ppt_stream = create_ppt_with_translated_text(uploaded_file, translated_slides)
            except Exception as e:
                st.error(f"PPT 파일을 생성하는 중 오류가 발생했습니다: {e}")
                return

        st.success("번역된 PPT 파일 생성 완료!")

        # 다운로드 버튼 제공
        st.download_button(
            label="영문 PPT 파일 다운로드",
            data=translated_ppt_stream,
            file_name=uploaded_file.name.split(".")[0]+"_en.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        st.balloons()

if __name__ == "__main__":
    main()